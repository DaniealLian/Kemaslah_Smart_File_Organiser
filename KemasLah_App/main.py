import sys
import threading
import logging
import os
import shutil
import re
from pathlib import Path

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QHBoxLayout,
    QStackedWidget, QVBoxLayout, QMessageBox, QDialog,
    QPushButton, QLabel, QProgressDialog, QTextEdit,
    QScrollArea, QFrame
)
from PyQt6.QtCore import pyqtSignal, Qt, QTimer, QThread

from auth.authentication_page import MainWindow as AuthWindow
from auth.server import app as flask_app

from src.gui.widgets.sidebar import Sidebar
from src.gui.widgets.topbar import TopBar
from src.gui.widgets.actionbar import ActionBar
from src.gui.widgets.loading_overlay import LoadingOverlay
from src.gui.views.home_view import HomeView
from src.gui.views.file_browser_view import FileBrowserView
from src.gui.views.archive_view import ArchiveView
from src.gui.views.statistics_view import StatisticsView, record_feature_use
from src.gui.views.settings_view import SettingsView

# --- Import the Share Dialog and File Sharing View ---
from src.gui.views.share_dialog import ShareFileDialog
from src.gui.views.file_sharing_view import FileSharingView
# -----------------------------------------------------

# import to add cnn model
from src.inference.classifier_worker import (
    CNNSearchWorker, IMAGE_EXTENSIONS, VIDEO_EXTENSIONS
)

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base_path, relative_path)

CNN_MODEL_PATH = resource_path(os.path.join("models", "trained", "best_model.pth"))

_MEDIA_EXTS = IMAGE_EXTENSIONS | VIDEO_EXTENSIONS


class DeepSearchWorker(QThread):
    progress = pyqtSignal(str)
    finished = pyqtSignal(dict)
    error = pyqtSignal(str)

    def __init__(self, text_files, ai_query):
        super().__init__()
        self.text_files = text_files
        self.ai_query = ai_query
        self.query_lower = ai_query.lower()
        self._is_cancelled = False

        self.supported_text_types = {
            'pdf', 'docx', 'txt', 'md', 'csv', 'py', 'json',
            'xlsx', 'xls', 'pptx', 'rtf', 'xml', 'html', 'htm',
            'yaml', 'yml', 'toml', 'ini', 'log', 'mhtml'
        }
        self.unsupported_types = {
            'zip', 'rar', '7z', 'tar', 'gz',
            'exe', 'dll', 'bin', 'iso', 'dmg',
            'db', 'sqlite', 'pkl', 'pyc'
        }

    def stop(self):
        self._is_cancelled = True

    def quick_content_extract(self, filepath, unsupported_files):
        ext = filepath.lower().split('.')[-1] if '.' in filepath else ''
        raw_text = ""

        if ext in self.unsupported_types:
            unsupported_files.append(os.path.basename(filepath))
            return ""

        try:
            if ext == 'pdf':
                import PyPDF2
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for i in range(min(5, len(reader.pages))):
                        if self._is_cancelled:
                            return ""
                        page_text = reader.pages[i].extract_text()
                        if page_text:
                            raw_text += page_text + " "

            elif ext == 'docx':
                import docx
                doc = docx.Document(filepath)
                for para in doc.paragraphs[:50]:
                    if self._is_cancelled:
                        return ""
                    raw_text += para.text + " "

            elif ext in ['xlsx', 'xls']:
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(max_row=100, values_only=True):
                        if self._is_cancelled:
                            wb.close()
                            return ""
                        raw_text += " ".join(str(c) for c in row if c is not None) + " "
                wb.close()

            elif ext == 'pptx':
                from pptx import Presentation as PptxPresentation
                prs = PptxPresentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if self._is_cancelled:
                            return ""
                        if hasattr(shape, "text"):
                            raw_text += shape.text + " "

            elif ext in [
                'txt', 'md', 'csv', 'py', 'json', 'rtf', 'xml',
                'html', 'htm', 'yaml', 'yml', 'toml', 'ini', 'log'
            ]:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    raw_text = f.read(10000)

            elif ext not in self.supported_text_types:
                unsupported_files.append(os.path.basename(filepath))
                return ""

        except Exception as read_err:
            print(f"Error reading file {filepath}: {read_err}")
            unsupported_files.append(os.path.basename(filepath))
            return ""

        return raw_text.lower()

    def run(self):
        try:
            unsupported_files = []
            docs = []
            valid_files = []

            total_files = len(self.text_files)

            for i, fp in enumerate(self.text_files, start=1):
                if self._is_cancelled:
                    self.finished.emit({
                        "cancelled": True,
                        "unsupported_files": unsupported_files,
                        "valid_files": [],
                        "found_matches": []
                    })
                    return

                self.progress.emit(f"Scanning {i}/{total_files}: {os.path.basename(fp)}")
                content = self.quick_content_extract(fp, unsupported_files)

                if content.strip():
                    docs.append(content)
                    valid_files.append(fp)

            found_matches = []

            if valid_files and not self._is_cancelled:
                for idx, content in enumerate(docs):
                    if self.query_lower in content:
                        found_matches.append((idx, 1.0, "Exact Keyword Match"))

                if not found_matches:
                    from sklearn.feature_extraction.text import TfidfVectorizer
                    from sklearn.metrics.pairwise import cosine_similarity

                    self.progress.emit("Running semantic AI matching...")

                    vectorizer = TfidfVectorizer(token_pattern=r'(?u)\b\w+\b')
                    tfidf_matrix = vectorizer.fit_transform([self.query_lower] + docs)
                    cosine_sim = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:]).flatten()

                    top_indices = cosine_sim.argsort()[-5:][::-1]
                    for idx in top_indices:
                        if cosine_sim[idx] > 0.001:
                            found_matches.append((idx, cosine_sim[idx], "Semantic AI Match"))

            self.finished.emit({
                "cancelled": False,
                "unsupported_files": unsupported_files,
                "valid_files": valid_files,
                "found_matches": found_matches
            })

        except Exception as e:
            self.error.emit(str(e))


class SmartOrganiseWorker(QThread):
    progress = pyqtSignal(str)
    finished = pyqtSignal(dict)
    error = pyqtSignal(str)

    def __init__(self, selected_paths, current_view_path, cnn_model_path):
        super().__init__()
        self.selected_paths = selected_paths
        self.current_view_path = current_view_path
        self.cnn_model_path = cnn_model_path
        self._is_cancelled = False

        self.supported_text_types = {
            'pdf', 'docx', 'txt', 'md', 'csv', 'py', 'json',
            'xlsx', 'xls', 'pptx', 'rtf', 'xml', 'html', 'htm',
            'yaml', 'yml', 'toml', 'ini', 'log'
        }
        self.unsupported_types = {
            'zip', 'rar', '7z', 'tar', 'gz',
            'exe', 'dll', 'bin', 'iso', 'dmg',
            'db', 'sqlite', 'pkl', 'pyc'
        }

    def stop(self):
        self._is_cancelled = True

    def _ensure_not_cancelled(self):
        if self._is_cancelled:
            raise InterruptedError("Smart organise cancelled.")

    def read_ultimate_precision_content(self, filepath, unsupported_files):
        self._ensure_not_cancelled()

        file_name = os.path.basename(filepath)
        ext = file_name.lower().split('.')[-1] if '.' in file_name else ''

        clean_name = re.sub(r'[^a-zA-Z\s]', ' ', file_name.replace('_', ' ').replace('-', ' '))
        content = (clean_name + " ") * 20

        if ext in self.unsupported_types:
            unsupported_files.append(file_name)
            return ""

        raw_text = ""
        try:
            if ext == 'pdf':
                import PyPDF2
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for i in range(min(6, len(reader.pages))):
                        self._ensure_not_cancelled()
                        page_text = reader.pages[i].extract_text()
                        if page_text:
                            if i == 0:
                                raw_text += (page_text[:2000] + " ") * 10
                            raw_text += page_text + " "

            elif ext == 'docx':
                import docx
                doc = docx.Document(filepath)
                for para in doc.paragraphs[:40]:
                    self._ensure_not_cancelled()
                    raw_text += para.text + " "

            elif ext in ['xlsx', 'xls']:
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(max_row=100, values_only=True):
                        self._ensure_not_cancelled()
                        raw_text += " ".join(str(c) for c in row if c is not None) + " "
                wb.close()

            elif ext == 'pptx':
                from pptx import Presentation as PptxPresentation
                prs = PptxPresentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        self._ensure_not_cancelled()
                        if hasattr(shape, "text"):
                            raw_text += shape.text + " "

            elif ext in [
                'txt', 'md', 'csv', 'py', 'json', 'rtf', 'xml',
                'html', 'htm', 'yaml', 'yml', 'toml', 'ini', 'log'
            ]:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    raw_text = f.read(15000)

            elif ext not in self.supported_text_types:
                unsupported_files.append(file_name)
                return ""

        except Exception as e:
            print(f"Extraction error for {file_name}: {e}")
            unsupported_files.append(file_name)
            return ""

        raw_text = raw_text.lower()
        words = re.findall(r'\b[a-z]{3,15}\b', raw_text)
        stemmed_words = [re.sub(r'(ing|tion|ment|ies|s)$', '', w) for w in words]
        raw_text = " ".join(stemmed_words)

        academic_noise = (
            r'\b(chapter|page|university|note|lecture|assignment|pdf|introduction|case|study|faculty|tarc|student|course|module|www|http|com|appendix|reference|conclusion|objective|assessment|rubric|tutorial|practical|guideline|data|result|table)\b'
        )
        raw_text = re.sub(academic_noise, '', raw_text)

        return content + re.sub(r'\s+', ' ', raw_text).strip()

    def _move_file_safely(self, source_path, dest_folder_path, move_history):
        self._ensure_not_cancelled()

        os.makedirs(dest_folder_path, exist_ok=True)

        file_name = os.path.basename(source_path)
        dest_file_path = os.path.join(dest_folder_path, file_name)

        base, ext_name = os.path.splitext(file_name)
        counter = 1
        while os.path.exists(dest_file_path):
            dest_file_path = os.path.join(dest_folder_path, f"{base}_copy{counter}{ext_name}")
            counter += 1

        shutil.move(source_path, dest_file_path)
        move_history.append((dest_file_path, source_path))
        return dest_file_path

    def _classify_and_move_media(self, image_paths, video_paths, dest_base, move_history):
        from src.inference.classifier import ImageClassifier, extract_keyframe

        results_msg = ""
        metrics_msg = ""
        temp_kfs = []
        classify_paths = list(image_paths)
        kf_to_video = {}

        keyframe_dir = os.path.join(os.path.expanduser("~"), ".kemaslah_temp", "keyframes")
        os.makedirs(keyframe_dir, exist_ok=True)

        for vp in video_paths:
            self._ensure_not_cancelled()
            self.progress.emit(f"Extracting keyframe: {os.path.basename(vp)}")
            kf = extract_keyframe(vp, output_dir="outputs/keyframes")
            if kf:
                classify_paths.append(kf)
                kf_to_video[kf] = vp
                temp_kfs.append(kf)
            else:
                results_msg += f"⚠️ Could not read keyframe from '{os.path.basename(vp)}' — skipped.\n"

        if not classify_paths:
            return results_msg, metrics_msg

        if not os.path.exists(self.cnn_model_path):
            results_msg += "❌ CNN model missing — images not sorted.\n"
            metrics_msg += "❌ CNN model missing — images not sorted.\n\n"
            return results_msg, metrics_msg

        self.progress.emit("Loading CNN model...")
        classifier = ImageClassifier(self.cnn_model_path)

        self.progress.emit(f"Classifying {len(classify_paths)} media file(s)...")
        cnn_results = classifier.classify_batch(classify_paths, batch_size=16)

        confident = 0
        fallback = 0

        for idx, result in enumerate(cnn_results, start=1):
            self._ensure_not_cancelled()

            clf_path = result["file_path"]
            category = result["category"]
            confidence = result["confidence"]
            accepted = result["accepted"]

            real_path = kf_to_video.get(clf_path, clf_path)
            file_name = os.path.basename(real_path)

            self.progress.emit(f"Moving media {idx}/{len(cnn_results)}: {file_name}")
            dest_folder = os.path.join(dest_base, category)
            self._move_file_safely(real_path, dest_folder, move_history)

            status = "✓" if accepted else "⚡"
            results_msg += f"{status} CNN: '{file_name}' ➡️ [{category}] ({confidence:.0%})\n"

            if accepted:
                confident += 1
            else:
                fallback += 1

        for kf in temp_kfs:
            try:
                os.remove(kf)
            except Exception:
                pass

        total = len(classify_paths)
        metrics_msg += (
            f"📊 CNN IMAGE CLASSIFICATION REPORT:\n"
            f"   • Backbone          : ResNet50 (Places365 + MS COCO)\n"
            f"   • Output categories : 10 KemasLah categories\n"
            f"   • Files processed   : {total}\n"
            f"   • High confidence   : {confident}  ({confident/total:.0%} of media)\n"
            f"   • Fallback category : {fallback}\n\n"
        )
        return results_msg, metrics_msg

    def run(self):
        try:
            folders_selected = [p for p in self.selected_paths if os.path.isdir(p)]
            files_selected = [p for p in self.selected_paths if os.path.isfile(p)]

            image_files = [f for f in files_selected if Path(f).suffix.lower() in IMAGE_EXTENSIONS]
            video_files = [f for f in files_selected if Path(f).suffix.lower() in VIDEO_EXTENSIONS]
            text_files = [f for f in files_selected if f not in image_files and f not in video_files]

            unsupported_organise_files = []
            valid_text_files = []

            for fp in text_files:
                fname = os.path.basename(fp)
                ext = fname.lower().split('.')[-1] if '.' in fname else ''
                if ext in self.unsupported_types:
                    unsupported_organise_files.append(fname)
                else:
                    valid_text_files.append(fp)

            text_files = valid_text_files

            move_history = []
            results_message = "🤖 AI Smart Organise Complete!\n\n"
            metrics_message = ""

            if len(files_selected) == 1 and not folders_selected and len(text_files) == 1:
                from sklearn.feature_extraction.text import TfidfVectorizer

                file_path = text_files[0]
                file_name = os.path.basename(file_path)

                self.progress.emit(f"Reading: {file_name}")
                text = self.read_ultimate_precision_content(file_path, unsupported_organise_files)

                if not text.strip():
                    self.finished.emit({
                        "cancelled": False,
                        "unsupported_files": unsupported_organise_files,
                        "error_message": (
                            f"'{file_name}' could not be read — it may have no text content "
                            "(e.g. image-only PDF).\n\nTry a different file format or ensure "
                            "the file contains readable text."
                        ),
                        "move_history": [],
                        "results_message": "",
                        "metrics_message": "",
                        "current_view_path": self.current_view_path
                    })
                    return

                self.progress.emit(f"Analysing keyword: {file_name}")
                vectorizer = TfidfVectorizer(stop_words='english', max_features=500, ngram_range=(1, 2))
                X = vectorizer.fit_transform([text])
                feature_names = vectorizer.get_feature_names_out()
                scores = X.toarray()[0]

                sorted_idx = scores.argsort()[::-1]
                target_folder = "Uncategorised"
                for idx in sorted_idx:
                    word = feature_names[idx].title()
                    if len(word) >= 4:
                        target_folder = word
                        break

                dest_folder_path = os.path.join(self.current_view_path, target_folder)
                self.progress.emit(f"Moving file to [{target_folder}]...")
                self._move_file_safely(file_path, dest_folder_path, move_history)

                results_message += f"✓ Sorted '{file_name}' ➡️ [{target_folder}]\n"
                metrics_message += (
                    f"📊 AI CONTENT ANALYSIS REPORT:\n"
                    f"   • Folder Name: {target_folder}\n"
                    f"   • Based On: Top keyword from file content\n"
                    f"   • Method: TF-IDF Keyword Extraction\n\n"
                )

            elif len(folders_selected) >= 2 and not files_selected:
                from sklearn.feature_extraction.text import TfidfVectorizer
                from sklearn.metrics.pairwise import cosine_similarity

                folder_profiles = []
                valid_folders = []

                self.progress.emit("Building folder profiles...")
                for i, folder in enumerate(folders_selected, start=1):
                    self._ensure_not_cancelled()
                    folder_name = os.path.basename(folder)
                    self.progress.emit(f"Profiling folder {i}/{len(folders_selected)}: {folder_name}")

                    combined_text = (folder_name + " ") * 10
                    for root, dirs, files in os.walk(folder):
                        for file in files:
                            self._ensure_not_cancelled()
                            combined_text += self.read_ultimate_precision_content(
                                os.path.join(root, file),
                                unsupported_organise_files
                            )
                            if len(combined_text) > 80000:
                                break

                    folder_profiles.append(combined_text)
                    valid_folders.append(folder)

                if len(folder_profiles) >= 2:
                    self.progress.emit("Calculating folder similarity...")
                    vectorizer = TfidfVectorizer(
                        stop_words='english',
                        max_features=2500,
                        ngram_range=(1, 2),
                        sublinear_tf=True
                    )

                    tfidf_matrix = vectorizer.fit_transform(folder_profiles)
                    sim_matrix = cosine_similarity(tfidf_matrix)

                    best_i, best_j, max_sim = 0, 1, -1
                    for i in range(len(valid_folders)):
                        for j in range(i + 1, len(valid_folders)):
                            if sim_matrix[i][j] > max_sim:
                                max_sim = sim_matrix[i][j]
                                best_i, best_j = i, j

                    if max_sim > -1:
                        folder_a, folder_b = valid_folders[best_i], valid_folders[best_j]
                        parent = folder_a if len(os.listdir(folder_a)) >= len(os.listdir(folder_b)) else folder_b
                        child = folder_b if parent == folder_a else folder_a

                        self.progress.emit(f"Merging folder [{os.path.basename(child)}]...")
                        target_path = os.path.join(parent, os.path.basename(child))
                        if os.path.exists(target_path):
                            target_path += "_merged"

                        shutil.move(child, target_path)
                        move_history.append((target_path, child))

                        results_message += (
                            f"📁 Force Merged Folders: [{os.path.basename(child)}] ➡️ "
                            f"[{os.path.basename(parent)}]\n"
                            f"   • Best Available Match Score: {max_sim:.1%}\n"
                        )

            elif folders_selected and files_selected:
                if image_files or video_files:
                    media_results, media_metrics = self._classify_and_move_media(
                        image_files, video_files, self.current_view_path, move_history
                    )
                    results_message += media_results
                    metrics_message += media_metrics

                if text_files:
                    import numpy as np
                    from sklearn.feature_extraction.text import TfidfVectorizer
                    from sklearn.ensemble import RandomForestClassifier
                    from sklearn.svm import SVC
                    from sklearn.metrics import precision_score

                    training_texts, training_labels = [], []
                    folder_paths_map = {}

                    self.progress.emit("Preparing training data from selected folders...")
                    for i, folder in enumerate(folders_selected, start=1):
                        self._ensure_not_cancelled()
                        folder_name = os.path.basename(folder)
                        folder_paths_map[folder_name] = folder

                        self.progress.emit(f"Reading folder {i}/{len(folders_selected)}: {folder_name}")
                        synthetic_data = (folder_name + " ") * 50
                        for _ in range(5):
                            training_texts.append(synthetic_data)
                            training_labels.append(folder_name)

                        for root, dirs, files in os.walk(folder):
                            for file in files:
                                self._ensure_not_cancelled()
                                text = self.read_ultimate_precision_content(
                                    os.path.join(root, file),
                                    unsupported_organise_files
                                )
                                if text.strip():
                                    training_texts.append(text)
                                    training_labels.append(folder_name)

                    unsorted_texts, valid_unsorted_files = [], []
                    self.progress.emit("Reading selected text files...")
                    for i, f in enumerate(text_files, start=1):
                        self._ensure_not_cancelled()
                        self.progress.emit(f"Reading file {i}/{len(text_files)}: {os.path.basename(f)}")
                        text = self.read_ultimate_precision_content(f, unsupported_organise_files)
                        if text.strip():
                            unsorted_texts.append(text)
                            valid_unsorted_files.append(f)

                    if valid_unsorted_files:
                        unique_labels = list(set(training_labels))

                        if len(unique_labels) < 2:
                            target_folder_name = unique_labels[0]
                            dest_folder_path = folder_paths_map[target_folder_name]

                            for i, file_path in enumerate(valid_unsorted_files, start=1):
                                self._ensure_not_cancelled()
                                self.progress.emit(
                                    f"Moving file {i}/{len(valid_unsorted_files)} to [{target_folder_name}]..."
                                )
                                file_name = os.path.basename(file_path)
                                self._move_file_safely(file_path, dest_folder_path, move_history)
                                results_message += (
                                    f"✓ Sorted '{file_name}' ➡️ [{target_folder_name}] "
                                    "(Single Target Direct Move)\n"
                                )

                            metrics_message += (
                                "📊 DIRECT SORTING:\n"
                                "   • Logic: Single Destination Selected\n"
                                "   • Coverage: 100% Automated Distribution\n\n"
                            )

                        else:
                            self.progress.emit("Training ensemble models...")
                            vectorizer = TfidfVectorizer(
                                stop_words='english',
                                max_features=5000,
                                ngram_range=(1, 3),
                                min_df=1,
                                max_df=0.80,
                                sublinear_tf=True
                            )
                            X_all = vectorizer.fit_transform(training_texts)

                            rf_model = RandomForestClassifier(
                                n_estimators=300,
                                criterion='entropy',
                                class_weight='balanced',
                                random_state=42
                            )
                            svm_model_sup = SVC(
                                kernel='linear',
                                C=3.0,
                                class_weight='balanced',
                                probability=True,
                                random_state=42
                            )

                            rf_model.fit(X_all, training_labels)
                            svm_model_sup.fit(X_all, training_labels)

                            rf_preds = rf_model.predict(X_all)
                            svm_preds = svm_model_sup.predict(X_all)

                            metrics_message += (
                                f"📊 ENSEMBLE CONSENSUS REPORT:\n"
                                f"   • RF Precision: {precision_score(training_labels, rf_preds, average='weighted', zero_division=0):.2%}\n"
                                f"   • SVM Precision: {precision_score(training_labels, svm_preds, average='weighted', zero_division=0):.2%}\n"
                                f"   • Logic: Soft Voting Probability Averaging\n\n"
                            )

                            X_unsorted = vectorizer.transform(unsorted_texts)
                            rf_probs = rf_model.predict_proba(X_unsorted)
                            svm_probs = svm_model_sup.predict_proba(X_unsorted)

                            rf_classes = rf_model.classes_
                            svm_classes = svm_model_sup.classes_

                            for i, file_path in enumerate(valid_unsorted_files, start=1):
                                self._ensure_not_cancelled()
                                self.progress.emit(
                                    f"Sorting text file {i}/{len(valid_unsorted_files)}: {os.path.basename(file_path)}"
                                )

                                rf_pred = rf_classes[np.argmax(rf_probs[i - 1])]
                                svm_pred = svm_classes[np.argmax(svm_probs[i - 1])]
                                rf_conf = np.max(rf_probs[i - 1])
                                svm_conf = np.max(svm_probs[i - 1])

                                if rf_pred == svm_pred:
                                    target_folder_name = rf_pred
                                    status_icon = "✓"
                                else:
                                    target_folder_name = rf_pred if rf_conf >= svm_conf else svm_pred
                                    status_icon = "⚡ (Tie-Breaker)"

                                dest_folder_path = folder_paths_map.get(target_folder_name)

                                if not dest_folder_path or not os.path.exists(dest_folder_path):
                                    if rf_pred == svm_pred:
                                        dest_folder_path = os.path.join(self.current_view_path, rf_pred)
                                    else:
                                        results_message += (
                                            f"⚠️ Skipped '{os.path.basename(file_path)}' ➡️ "
                                            f"[{target_folder_name}] (Hallucinated Folder Target)\n"
                                        )
                                        continue

                                self._move_file_safely(file_path, dest_folder_path, move_history)
                                results_message += (
                                    f"{status_icon} Sorted '{os.path.basename(file_path)}' ➡️ "
                                    f"[{target_folder_name}]\n"
                                )

            elif (len(files_selected) >= 2 and not folders_selected) or (folders_selected and not files_selected):
                pool = list(files_selected)
                for folder in folders_selected:
                    for root, dirs, files in os.walk(folder):
                        for file in files:
                            pool.append(os.path.join(root, file))

                pool_images = [f for f in pool if Path(f).suffix.lower() in IMAGE_EXTENSIONS]
                pool_videos = [f for f in pool if Path(f).suffix.lower() in VIDEO_EXTENSIONS]
                pool_text = [f for f in pool if f not in pool_images and f not in pool_videos]

                if pool_images or pool_videos:
                    media_results, media_metrics = self._classify_and_move_media(
                        pool_images, pool_videos, self.current_view_path, move_history
                    )
                    results_message += media_results
                    metrics_message += media_metrics

                if pool_text:
                    import numpy as np
                    from sklearn.feature_extraction.text import TfidfVectorizer
                    from sklearn.cluster import KMeans
                    from sklearn.ensemble import RandomForestClassifier
                    from sklearn.svm import SVC

                    file_contents, valid_files = [], []

                    self.progress.emit("Reading files for clustering...")
                    for i, f in enumerate(pool_text, start=1):
                        self._ensure_not_cancelled()
                        self.progress.emit(f"Reading file {i}/{len(pool_text)}: {os.path.basename(f)}")
                        text = self.read_ultimate_precision_content(f, unsupported_organise_files)
                        if text.strip():
                            file_contents.append(text)
                            valid_files.append(f)

                    if len(valid_files) >= 2:
                        self.progress.emit("Building clusters...")
                        vectorizer = TfidfVectorizer(
                            stop_words='english',
                            max_features=3000,
                            ngram_range=(1, 2),
                            min_df=1,
                            max_df=0.90,
                            sublinear_tf=True
                        )
                        X = vectorizer.fit_transform(file_contents)
                        feature_names = vectorizer.get_feature_names_out()

                        n_clusters = max(2, len(valid_files) // 3)
                        if n_clusters >= len(valid_files):
                            n_clusters = max(2, len(valid_files) - 1)
                        n_clusters = min(n_clusters, len(valid_files))

                        kmeans = KMeans(
                            n_clusters=n_clusters,
                            random_state=42,
                            n_init=30,
                            max_iter=1000
                        )
                        kmeans.fit(X)

                        cluster_labels = kmeans.labels_
                        unique_clusters = len(set(cluster_labels))

                        if unique_clusters < 2:
                            self.finished.emit({
                                "cancelled": False,
                                "unsupported_files": unsupported_organise_files,
                                "error_message": (
                                    "These text files are too similar for the AI to separate into different groups.\n\n"
                                    "Try selecting files with more varied content, or manually sort them into folders first."
                                ),
                                "move_history": move_history,
                                "results_message": results_message,
                                "metrics_message": metrics_message,
                                "current_view_path": self.current_view_path
                            })
                            return

                        rf_model_unsup = RandomForestClassifier(
                            n_estimators=300,
                            criterion='entropy',
                            class_weight='balanced',
                            random_state=42
                        )
                        svm_model_unsup = SVC(
                            kernel='linear',
                            C=2.0,
                            class_weight='balanced',
                            probability=True,
                            random_state=42
                        )

                        self.progress.emit("Training cluster boundary models...")
                        rf_model_unsup.fit(X, cluster_labels)
                        svm_model_unsup.fit(X, cluster_labels)

                        metrics_message += (
                            "📊 ENTERPRISE CLUSTER REPORT:\n"
                            "   • Logic: Dual-Model Probability Averaging\n"
                            "   • Status: 100% Boundary Assignment\n\n"
                        )

                        rf_probs = rf_model_unsup.predict_proba(X)
                        svm_probs = svm_model_unsup.predict_proba(X)
                        classes = svm_model_unsup.classes_

                        for i, file_path in enumerate(valid_files, start=1):
                            self._ensure_not_cancelled()
                            self.progress.emit(
                                f"Grouping file {i}/{len(valid_files)}: {os.path.basename(file_path)}"
                            )

                            avg_probs = (rf_probs[i - 1] + svm_probs[i - 1]) / 2.0
                            assigned_label = classes[np.argmax(avg_probs)]

                            top_word_idx = kmeans.cluster_centers_[assigned_label].argsort()[-1]
                            target_folder = feature_names[top_word_idx].title()
                            if len(target_folder) < 3:
                                target_folder = f"Subject_Area_{assigned_label + 1}"

                            dest_folder_path = os.path.join(self.current_view_path, target_folder)
                            self._move_file_safely(file_path, dest_folder_path, move_history)
                            results_message += (
                                f"✓ Smart Grouped '{os.path.basename(file_path)}' ➡️ "
                                f"[{target_folder}]\n"
                            )

            else:
                self.finished.emit({
                    "cancelled": False,
                    "unsupported_files": unsupported_organise_files,
                    "error_message": "Please select files or folders to organize.",
                    "move_history": [],
                    "results_message": "",
                    "metrics_message": "",
                    "current_view_path": self.current_view_path
                })
                return

            self.finished.emit({
                "cancelled": False,
                "unsupported_files": unsupported_organise_files,
                "error_message": "",
                "move_history": move_history,
                "results_message": results_message,
                "metrics_message": metrics_message,
                "current_view_path": self.current_view_path
            })

        except InterruptedError:
            self.finished.emit({
                "cancelled": True,
                "unsupported_files": [],
                "error_message": "",
                "move_history": [],
                "results_message": "",
                "metrics_message": "",
                "current_view_path": self.current_view_path
            })
        except Exception as e:
            self.error.emit(str(e))


class SmartFileManager(QMainWindow):
    """Enterprise File Manager: Absolute Keyword Search & Anti-Silent-Fail Logging"""
    logout_requested = pyqtSignal()

    def __init__(self, user_data=None, auth_window=None):
        super().__init__()
        self.user_data = user_data or {"username": "Guest", "email": "guest@local", "initials": "G"}
        self.auth_window = auth_window
        self.setWindowTitle("Kemaslah File Manager")
        self.resize(1200, 800)
        self.settings_overlay = None
        self.overlay = None

        self.search_timer = QTimer()
        self.search_timer.setSingleShot(True)
        self.search_timer.timeout.connect(self.perform_search)
        self.current_search_query = ""

        self._cnn_search_worker = None
        self._deep_search_worker = None
        self._smart_organise_worker = None

        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QHBoxLayout(central_widget)
        main_layout.setContentsMargins(0, 0, 0, 0)
        main_layout.setSpacing(0)

        self.sidebar = Sidebar(user_data=self.user_data)
        self.sidebar.navigation_changed.connect(self.switch_view)
        self.sidebar.logout_requested.connect(self.handle_logout)
        main_layout.addWidget(self.sidebar)

        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(0)

        self.top_bar = TopBar()
        self.top_bar.path_changed.connect(self.on_topbar_nav)
        self.top_bar.search_query_changed.connect(self.handle_search)
        self.top_bar.refresh_clicked.connect(self.handle_refresh)
        right_layout.addWidget(self.top_bar)

        self.stack = QStackedWidget()
        self.home_view = HomeView()
        self.files_view = FileBrowserView()
        self.archive_view = ArchiveView()
        self.statistics_view = StatisticsView()
        self.sharing_view = FileSharingView(self.user_data.get("email"))

        self.action_bar = ActionBar(True, "Smart Organise")
        self.action_bar.action_clicked.connect(self.handle_action_bar)
        self.action_bar.smart_organise_clicked.connect(self.handle_smart_organise)

        self.action_bar.nav_back_clicked.connect(self.files_view.go_back)
        self.action_bar.nav_forward_clicked.connect(self.files_view.go_forward)

        right_layout.addWidget(self.action_bar)

        self.files_view.path_changed.connect(self.top_bar.update_breadcrumbs)
        self.home_view.folder_opened.connect(self.on_home_folder_opened)
        self.files_view.file_table.share_requested.connect(self.open_share_dialog)

        self.stack.addWidget(self.home_view)
        self.stack.addWidget(self.files_view)
        self.stack.addWidget(self.archive_view)
        self.stack.addWidget(self.statistics_view)
        self.stack.addWidget(self.sharing_view)

        right_layout.addWidget(self.stack)
        main_layout.addWidget(right_panel)

        self.overlay = LoadingOverlay(self)
        self.overlay.resize(self.size())

        self.switch_view("home")

    # ── Overlay helpers ──────────────────────────────────────────────────────

    def _show_overlay(self, message):
        if self.overlay:
            self.overlay.show_message(message)

    def _update_overlay(self, message):
        if self.overlay and self.overlay.isVisible():
            self.overlay.label.setText(message)

    def _hide_overlay(self):
        if self.overlay:
            self.overlay.hide_overlay()

    def open_share_dialog(self, file_path):
        is_guest = (self.user_data.get("username") == "Guest User")
        dialog = ShareFileDialog(
            file_path=file_path,
            current_user_email=self.user_data.get("email"),
            is_guest=is_guest,
            parent=self
        )
        dialog.exec()

    def handle_logout(self):
        reply = QMessageBox.question(
            self,
            'Logout',
            'Are you sure you want to logout?',
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        if reply == QMessageBox.StandardButton.Yes:
            self.logout_requested.emit()

    def handle_action_bar(self, action_name):
        current_widget = self.stack.currentWidget()
        if current_widget == self.files_view:
            self.files_view.file_table.perform_action(action_name.lower())

    # -------------------------------------------------------------------------
    # --- ABSOLUTE SEARCH ENHANCEMENTS ---
    # -------------------------------------------------------------------------
    def handle_search(self, query):
        self.current_search_query = query
        if self.stack.currentWidget() != self.files_view:
            self.switch_view("files")
        self.files_view.file_table.filter_files(query)
        self.search_timer.start(800)

    def perform_search(self):
        query = self.current_search_query.strip()
        if len(query) < 3:
            return
        self.perform_ai_deep_search(query)

    def perform_ai_deep_search(self, ai_query):
        current_path = self.files_view.current_path
        if not os.path.exists(current_path):
            return

        all_files = [
            os.path.join(current_path, f)
            for f in os.listdir(current_path)
            if os.path.isfile(os.path.join(current_path, f))
        ]

        if not all_files:
            QMessageBox.information(
                self,
                "Content Scanner",
                f"There are no files inside the folder '{os.path.basename(current_path)}' to scan.\n\n"
                "Please navigate to the exact folder where your file is saved!"
            )
            return

        text_files = [f for f in all_files if Path(f).suffix.lower() not in _MEDIA_EXTS]
        media_files = [f for f in all_files if Path(f).suffix.lower() in _MEDIA_EXTS]

        if text_files:
            if self._deep_search_worker and self._deep_search_worker.isRunning():
                self._deep_search_worker.stop()
                self._deep_search_worker.wait()

            self._show_overlay("Searching inside files...")

            self._deep_search_worker = DeepSearchWorker(text_files, ai_query)
            self._deep_search_worker.progress.connect(self._on_deep_search_progress)
            self._deep_search_worker.finished.connect(lambda result: self._on_deep_search_finished(result, ai_query))
            self._deep_search_worker.error.connect(self._on_deep_search_error)
            self._deep_search_worker.start()

        if media_files:
            if not os.path.exists(CNN_MODEL_PATH):
                print("[CNN Search] best_model.pth not found — skipping image search.")
                return

            if self._cnn_search_worker and self._cnn_search_worker.isRunning():
                self._cnn_search_worker.stop()
                self._cnn_search_worker.wait()

            QMessageBox.information(
                self,
                "Image Search Started",
                f"🖼️  Also scanning {len(media_files)} image/video file(s) by "
                f"visual content.\n\nMatching results will appear in the "
                f"file list below as they are found."
            )

            self._cnn_search_worker = CNNSearchWorker(
                query=ai_query,
                search_path=current_path,
                model_path=CNN_MODEL_PATH,
            )
            self._cnn_search_worker.match_found.connect(self._on_cnn_image_match)
            self._cnn_search_worker.search_finished.connect(self._on_cnn_search_done)
            self._cnn_search_worker.error_occurred.connect(self._on_cnn_search_error)
            self._cnn_search_worker.start()

    def _cancel_deep_search(self):
        if self._deep_search_worker and self._deep_search_worker.isRunning():
            self._deep_search_worker.stop()
        self._hide_overlay()

    def _on_deep_search_progress(self, message):
        self._update_overlay(message)

    def _on_deep_search_finished(self, result, ai_query):
        self._hide_overlay()

        if result.get("cancelled"):
            QMessageBox.information(self, "Deep Search", "Deep search was cancelled.")
            return

        unsupported_files = result.get("unsupported_files", [])
        valid_files = result.get("valid_files", [])
        found_matches = result.get("found_matches", [])

        if unsupported_files:
            QMessageBox.warning(
                self,
                "Unsupported File Types",
                f"The following {len(unsupported_files)} file(s) could not be read (no text content):\n\n"
                + "\n".join(f"• {f}" for f in unsupported_files)
                + "\n\nSupported types: PDF, DOCX, XLSX, PPTX, TXT, CSV, MD, PY, JSON, HTML, XML, YAML, LOG, etc."
            )

        if valid_files and found_matches:
            result_msg = f"🔍 Inside-File Matches for: '{ai_query}'\n\n"
            for idx, score, match_type in found_matches:
                file_name = os.path.basename(valid_files[idx])
                score_txt = "100% (Exact Match)" if match_type == "Exact Keyword Match" else f"{score:.1%}"
                result_msg += (
                    f"📄 {file_name}\n"
                    f"   • Match Type: {match_type}\n"
                    f"   • Relevance: {score_txt}\n\n"
                )

            dialog = QDialog(self)
            dialog.setWindowTitle("File Content Search Results")
            dialog.resize(600, 400)
            layout = QVBoxLayout(dialog)

            instruction_lbl = QLabel(
                f"<b>Top matches based on inside document text for '{ai_query}':</b><br>"
                "<i>(Note: Files that don't have this word in their title are hidden behind this box)</i>"
            )
            layout.addWidget(instruction_lbl)

            text_edit = QTextEdit(dialog)
            text_edit.setReadOnly(True)
            text_edit.setText(result_msg)
            layout.addWidget(text_edit)

            btn = QPushButton("Close Search Results")
            btn.clicked.connect(dialog.accept)
            layout.addWidget(btn)
            dialog.exec()

    def _on_deep_search_error(self, message):
        self._hide_overlay()

        if "empty vocabulary" in str(message).lower():
            QMessageBox.warning(self, "AI Search Alert", "Search query contains invalid characters.")
        else:
            QMessageBox.critical(
                self,
                "System Error",
                f"The AI encountered an error while scanning the text:\n\n{message}\n\nCheck terminal for details."
            )
            print(f"Deep Search Error: {message}")

    # -------------------------------------------------------------------------
    # --- SMART ORGANISE ---
    # -------------------------------------------------------------------------
    def handle_smart_organise(self):
        current_widget = self.stack.currentWidget()

        if current_widget == self.archive_view:
            self.archive_view.open_date_dialog()
            return

        if current_widget != self.files_view:
            return

        selected_paths = self.files_view.file_table.get_selected_files()
        if not selected_paths:
            QMessageBox.warning(self, "No Selection", "Please select folders or files to organize.")
            return

        if self._smart_organise_worker and self._smart_organise_worker.isRunning():
            QMessageBox.information(self, "Smart Organise", "Smart Organise is already running.")
            return

        self._show_overlay("Organizing files...")

        self._smart_organise_worker = SmartOrganiseWorker(
            selected_paths=selected_paths,
            current_view_path=self.files_view.current_path,
            cnn_model_path=CNN_MODEL_PATH
        )
        self._smart_organise_worker.progress.connect(self._on_smart_organise_progress)
        self._smart_organise_worker.finished.connect(self._on_smart_organise_finished)
        self._smart_organise_worker.error.connect(self._on_smart_organise_error)
        self._smart_organise_worker.start()

    def _cancel_smart_organise(self):
        if self._smart_organise_worker and self._smart_organise_worker.isRunning():
            self._smart_organise_worker.stop()
        self._hide_overlay()

    def _on_smart_organise_progress(self, message):
        self._update_overlay(message)

    def _on_smart_organise_finished(self, result):
        self._hide_overlay()

        if result.get("cancelled"):
            QMessageBox.information(self, "Smart Organise", "Smart Organise was cancelled.")
            return

        unsupported_files = result.get("unsupported_files", [])
        error_message = result.get("error_message", "")
        move_history = result.get("move_history", [])
        results_message = result.get("results_message", "")
        metrics_message = result.get("metrics_message", "")
        current_view_path = result.get("current_view_path", self.files_view.current_path)

        if unsupported_files:
            QMessageBox.warning(
                self,
                "Unsupported File Types",
                f"The following {len(unsupported_files)} file(s) cannot be analysed (no text content) and will be skipped:\n\n"
                + "\n".join(f"• {f}" for f in unsupported_files)
                + "\n\nSupported types: PDF, DOCX, XLSX, PPTX, TXT, CSV, MD, PY, JSON, HTML, XML, YAML, LOG, etc."
            )

        if error_message:
            QMessageBox.warning(self, "Smart Organise", error_message)

        if move_history:
            self.handle_satisfaction_check(
                results_message,
                metrics_message,
                move_history,
                current_view_path
            )
        else:
            if hasattr(self.files_view, 'file_table'):
                self.files_view.file_table.load_files(current_view_path)

    def _on_smart_organise_error(self, message):
        self._hide_overlay()
        QMessageBox.critical(self, "Semantic Error", f"Precision pipeline failed: {message}")

    def handle_satisfaction_check(self, results, metrics, history, current_path):
        if history:
            record_feature_use("organise")
        if hasattr(self.files_view, 'file_table'):
            self.files_view.file_table.load_files(current_path)

        if history:
            dialog = QDialog(self)
            dialog.setWindowTitle("Smart Organise Results")
            dialog.setMinimumSize(780, 600)
            dialog.setStyleSheet("""
                QDialog {
                    background-color: #0F172A;
                    color: #F1F5F9;
                }
            """)

            root = QVBoxLayout(dialog)
            root.setContentsMargins(0, 0, 0, 0)
            root.setSpacing(0)

            header = QWidget()
            header.setFixedHeight(64)
            header.setStyleSheet("background-color: #1E293B; border-bottom: 1px solid #334155;")
            h_layout = QHBoxLayout(header)
            h_layout.setContentsMargins(24, 0, 24, 0)

            icon_lbl = QLabel("🤖")
            icon_lbl.setStyleSheet("font-size: 22px;")
            title_lbl = QLabel("AI Smart Organise — Results")
            title_lbl.setStyleSheet("font-size: 16px; font-weight: bold; color: #F1F5F9;")
            sub_lbl = QLabel(f"{len(history)} file(s) sorted")
            sub_lbl.setStyleSheet("font-size: 12px; color: #64748B; margin-left: 8px; margin-top: 3px;")

            h_layout.addWidget(icon_lbl)
            h_layout.addSpacing(10)
            h_layout.addWidget(title_lbl)
            h_layout.addWidget(sub_lbl)
            h_layout.addStretch()
            root.addWidget(header)

            body = QWidget()
            body.setStyleSheet("background-color: #0F172A;")
            body_layout = QHBoxLayout(body)
            body_layout.setContentsMargins(20, 20, 20, 20)
            body_layout.setSpacing(16)

            metrics_panel = QFrame()
            metrics_panel.setFixedWidth(220)
            metrics_panel.setStyleSheet("""
                QFrame {
                    background-color: #1E293B;
                    border-radius: 10px;
                    padding: 4px;
                }
            """)
            m_layout = QVBoxLayout(metrics_panel)
            m_layout.setContentsMargins(16, 16, 16, 16)
            m_layout.setSpacing(10)

            metrics_title = QLabel("📊 AI Report")
            metrics_title.setStyleSheet("font-size: 13px; font-weight: bold; color: #60A5FA;")
            m_layout.addWidget(metrics_title)

            for line in metrics.strip().split("\n"):
                line = line.strip().lstrip("•").strip()
                if not line or "REPORT" in line or "CLUSTER" in line or "ENSEMBLE" in line or "BBC" in line or "CONTENT" in line:
                    continue
                lbl = QLabel(line)
                lbl.setWordWrap(True)
                lbl.setStyleSheet("font-size: 11px; color: #94A3B8; padding: 2px 0;")
                m_layout.addWidget(lbl)

            m_layout.addStretch()

            q_lbl = QLabel("Are you satisfied\nwith this organisation?")
            q_lbl.setAlignment(Qt.AlignmentFlag.AlignCenter)
            q_lbl.setStyleSheet("""
                font-size: 12px; color: #CBD5E1;
                background-color: #0F172A;
                border-radius: 6px;
                padding: 8px;
            """)
            m_layout.addWidget(q_lbl)

            body_layout.addWidget(metrics_panel)

            right_panel = QWidget()
            right_panel.setStyleSheet("background-color: transparent;")
            r_layout = QVBoxLayout(right_panel)
            r_layout.setContentsMargins(0, 0, 0, 0)
            r_layout.setSpacing(8)

            results_title = QLabel("Sorted Files")
            results_title.setStyleSheet("font-size: 13px; font-weight: bold; color: #F1F5F9; margin-bottom: 4px;")
            r_layout.addWidget(results_title)

            scroll = QScrollArea()
            scroll.setWidgetResizable(True)
            scroll.setStyleSheet("""
                QScrollArea { border: none; background-color: transparent; }
                QScrollBar:vertical { background: #1E293B; width: 6px; border-radius: 3px; }
                QScrollBar::handle:vertical { background: #475569; border-radius: 3px; }
            """)

            scroll_widget = QWidget()
            scroll_widget.setStyleSheet("background-color: transparent;")
            scroll_layout = QVBoxLayout(scroll_widget)
            scroll_layout.setContentsMargins(0, 0, 0, 0)
            scroll_layout.setSpacing(6)

            for line in results.strip().split("\n"):
                line = line.strip()
                if not line or "AI Smart Organise" in line:
                    continue

                card = QFrame()
                card.setStyleSheet("""
                    QFrame {
                        background-color: #1E293B;
                        border-radius: 8px;
                        border: 1px solid #273549;
                    }
                """)
                card_layout = QHBoxLayout(card)
                card_layout.setContentsMargins(12, 10, 12, 10)
                card_layout.setSpacing(10)

                if line.startswith("✓"):
                    status_icon = QLabel("✓")
                    status_icon.setStyleSheet("color: #22C55E; font-size: 14px; font-weight: bold;")
                elif line.startswith("⚠"):
                    status_icon = QLabel("⚠")
                    status_icon.setStyleSheet("color: #F59E0B; font-size: 14px;")
                else:
                    status_icon = QLabel("→")
                    status_icon.setStyleSheet("color: #60A5FA; font-size: 14px;")

                status_icon.setFixedWidth(20)
                card_layout.addWidget(status_icon)

                text = line.lstrip("✓⚡⚠ ").strip()
                if "➡️" in text:
                    parts = text.split("➡️")
                    fname = parts[0].replace("Sorted", "").replace("Smart Grouped", "").replace("'", "").strip()
                    folder = parts[1].strip().strip("[]") if len(parts) > 1 else ""

                    file_lbl = QLabel(fname)
                    file_lbl.setStyleSheet("font-size: 12px; color: #E2E8F0;")
                    file_lbl.setWordWrap(True)
                    card_layout.addWidget(file_lbl, 1)

                    arrow = QLabel("→")
                    arrow.setStyleSheet("color: #475569; font-size: 13px;")
                    card_layout.addWidget(arrow)

                    folder_badge = QLabel(f"📁 {folder}")
                    folder_badge.setStyleSheet("""
                        font-size: 11px; color: #60A5FA;
                        background-color: #172554;
                        border-radius: 4px;
                        padding: 2px 8px;
                    """)
                    folder_badge.setFixedHeight(22)
                    card_layout.addWidget(folder_badge)
                else:
                    text_lbl = QLabel(text)
                    text_lbl.setStyleSheet("font-size: 12px; color: #94A3B8;")
                    text_lbl.setWordWrap(True)
                    card_layout.addWidget(text_lbl, 1)

                scroll_layout.addWidget(card)

            scroll_layout.addStretch()
            scroll.setWidget(scroll_widget)
            r_layout.addWidget(scroll)
            body_layout.addWidget(right_panel, 1)
            root.addWidget(body, 1)

            footer = QWidget()
            footer.setFixedHeight(68)
            footer.setStyleSheet("background-color: #1E293B; border-top: 1px solid #334155;")
            f_layout = QHBoxLayout(footer)
            f_layout.setContentsMargins(24, 0, 24, 0)

            hint = QLabel("💡 Undo will restore all files to their original locations.")
            hint.setStyleSheet("font-size: 11px; color: #64748B;")
            f_layout.addWidget(hint)
            f_layout.addStretch()

            undo_btn = QPushButton("↩  Undo Changes")
            undo_btn.setFixedSize(150, 38)
            undo_btn.setStyleSheet("""
                QPushButton {
                    background-color: #1E293B;
                    color: #94A3B8;
                    border: 1px solid #475569;
                    border-radius: 8px;
                    font-size: 13px;
                    font-weight: bold;
                }
                QPushButton:hover {
                    background-color: #334155;
                    color: #F1F5F9;
                }
            """)

            keep_btn = QPushButton("✔  Keep Changes")
            keep_btn.setFixedSize(150, 38)
            keep_btn.setStyleSheet("""
                QPushButton {
                    background-color: #2563EB;
                    color: white;
                    border: none;
                    border-radius: 8px;
                    font-size: 13px;
                    font-weight: bold;
                }
                QPushButton:hover {
                    background-color: #1D4ED8;
                }
            """)

            undo_btn.clicked.connect(dialog.reject)
            keep_btn.clicked.connect(dialog.accept)

            f_layout.addWidget(undo_btn)
            f_layout.addSpacing(10)
            f_layout.addWidget(keep_btn)
            root.addWidget(footer)

            if dialog.exec() == QDialog.DialogCode.Rejected:
                undo_progress = QProgressDialog("Restoring originals...", None, 0, len(history), self)
                undo_progress.setMinimumDuration(0)
                undo_progress.setWindowModality(Qt.WindowModality.WindowModal)
                for i, (new_pos, old_pos) in enumerate(reversed(history)):
                    if os.path.exists(new_pos):
                        shutil.move(new_pos, old_pos)
                        try:
                            if not os.listdir(os.path.dirname(new_pos)):
                                os.rmdir(os.path.dirname(new_pos))
                        except Exception:
                            pass
                    undo_progress.setValue(i + 1)
                    QApplication.processEvents()
                undo_progress.close()
                QMessageBox.information(self, "Restored", "All files and folders have been returned to their original positions.")
                if hasattr(self.files_view, 'file_table'):
                    self.files_view.file_table.load_files(current_path)

    def switch_view(self, identifier):
        lang_code = self.user_data.get('language_code', 'en')

        if identifier == "settings":
            self.show_settings_overlay()
        elif identifier in ["sharing", "file_sharing"]:
            self.sharing_view.load_shared_files()
            self.stack.setCurrentWidget(self.sharing_view)
            self.top_bar.update_breadcrumbs("File Sharing")
            self.sidebar.set_active(identifier)
        else:
            self.sidebar.set_active(identifier)
            if identifier == "home":
                self.stack.setCurrentWidget(self.home_view)
                self.top_bar.update_breadcrumbs("Home")
            elif identifier == "files":
                self.stack.setCurrentWidget(self.files_view)
                self.top_bar.update_breadcrumbs(self.files_view.current_path)
                self.action_bar.set_smart_mode("organise", lang_code)
            elif identifier == "archive":
                self.stack.setCurrentWidget(self.archive_view)
                self.top_bar.update_breadcrumbs("Smart Archive")
                self.action_bar.set_smart_mode("archive", lang_code)
            elif identifier == "statistics":
                self.stack.setCurrentWidget(self.statistics_view)
                self.top_bar.update_breadcrumbs("Statistics")

    def on_topbar_nav(self, path):
        if self.stack.currentWidget() == self.files_view:
            self.files_view.on_breadcrumb_clicked(path)

    def show_settings_overlay(self):
        if not self.settings_overlay:
            self.settings_overlay = SettingsView(self.user_data)
            self.settings_overlay.setParent(self)
            self.settings_overlay.closed.connect(self.hide_settings_overlay)
            self.settings_overlay.logout_requested.connect(self.handle_logout)
        self.settings_overlay.resize(self.size())
        self.settings_overlay.show()
        self.settings_overlay.raise_()

    def hide_settings_overlay(self):
        if self.settings_overlay:
            self.settings_overlay.hide()

    def resizeEvent(self, event):
        if self.settings_overlay and self.settings_overlay.isVisible():
            self.settings_overlay.resize(self.size())

        if self.overlay:
            self.overlay.resize(self.size())

        super().resizeEvent(event)

    def on_home_folder_opened(self, path):
        self.files_view.navigate_to(path)
        self.switch_view("files")
        self.sidebar.set_active("files")

    def handle_refresh(self):
        current = self.stack.currentWidget()
        if current == self.files_view:
            self.files_view.file_table.load_files(self.files_view.current_path)
        elif current == self.home_view:
            self.home_view.load_recent_files() if hasattr(self.home_view, 'load_recent_files') else None
            self.home_view.repaint()
        elif current == self.statistics_view:
            self.statistics_view._refresh() if hasattr(self.statistics_view, '_refresh') else None
        elif current == self.archive_view:
            self.archive_view.file_table.load_files(self.archive_view.file_table.current_path) \
                if hasattr(self.archive_view, 'file_table') and self.archive_view.file_table.current_path else None
        elif current == self.sharing_view:
            self.sharing_view.load_shared_files() if hasattr(self.sharing_view, 'load_shared_files') else None

    def mousePressEvent(self, event):
        if self.stack.currentWidget() == self.files_view:
            if event.button() == Qt.MouseButton.BackButton:
                self.files_view.go_back()
                event.accept()
                return
            elif event.button() == Qt.MouseButton.ForwardButton:
                self.files_view.go_forward()
                event.accept()
                return
        super().mousePressEvent(event)

    def _on_cnn_image_match(self, name: str, full_path: str, category: str, confidence: float):
        self.files_view.file_table._add_search_row(name, full_path, is_dir=False)

    def _on_cnn_search_done(self, scanned: int, matched: int):
        print(f"[CNN Search] Done — scanned {scanned} media files, matched {matched}.")

    def _on_cnn_search_error(self, message: str):
        print(f"[CNN Search] Error (non-fatal): {message}")


class KemaslahApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Kemaslah - Smart File Manager")
        self.stack = QStackedWidget()
        self.setCentralWidget(self.stack)
        self.auth_window = AuthWindow()
        self.stack.addWidget(self.auth_window)
        self.file_manager = None
        login_page = self.auth_window.stack.widget(0)

        if hasattr(login_page, 'login_successful'):
            login_page.login_successful.connect(self.on_login_success)
        if hasattr(login_page, 'skip_login_clicked'):
            login_page.skip_login_clicked.connect(self.on_skip_login)

        self.setMinimumSize(0, 0)
        self.setMaximumSize(16777215, 16777215)
        self.resize(1000, 750)

    def on_login_success(self, user_data):
        self._destroy_file_manager()
        self.file_manager = SmartFileManager(user_data, auth_window=self.auth_window)
        self.file_manager.logout_requested.connect(self.on_logout)
        self.stack.addWidget(self.file_manager)
        self.stack.setCurrentWidget(self.file_manager)
        self.file_manager.switch_view("home")
        self.setMinimumSize(1200, 800)
        self.setMaximumSize(16777215, 16777215)
        self.showMaximized()

    def on_skip_login(self):
        guest_data = {
            "username": "Guest User",
            "email": "guest@local",
            "initials": "GU",
            "display_name": "Guest User"
        }
        if self.file_manager is None:
            self.file_manager = SmartFileManager(guest_data, auth_window=None)
            self.file_manager.logout_requested.connect(self.on_guest_exit)
            self.stack.addWidget(self.file_manager)
        self.stack.setCurrentWidget(self.file_manager)
        self.file_manager.switch_view("home")
        self.setMinimumSize(1200, 800)
        self.setMaximumSize(16777215, 16777215)
        self.showMaximized()

    def on_guest_exit(self):
        reply = QMessageBox.question(
            self,
            'Logout',
            'Exit guest mode and return to login?',
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        if reply == QMessageBox.StandardButton.Yes:
            self._destroy_file_manager()
            self.stack.setCurrentWidget(self.auth_window)
            self.auth_window.stack.setCurrentIndex(0)
            self.setMinimumSize(0, 0)
            self.setMaximumSize(16777215, 16777215)
            self.resize(1000, 750)
            self.showNormal()

    def on_logout(self):
        self._destroy_file_manager()
        login_page = self.auth_window.stack.widget(0)
        if hasattr(login_page, 'clear_fields'):
            login_page.clear_fields()
        self.stack.setCurrentWidget(self.auth_window)
        self.setMinimumSize(0, 0)
        self.setMaximumSize(16777215, 16777215)
        self.resize(1000, 750)
        self.showNormal()
        self.auth_window.stack.setCurrentIndex(0)

    def _destroy_file_manager(self):
        if self.file_manager is not None:
            self.stack.removeWidget(self.file_manager)
            self.file_manager.deleteLater()
            self.file_manager = None

    def update_all_pages(self, lang_code):
        if hasattr(self.auth_window, 'update_all_pages'):
            self.auth_window.update_all_pages(lang_code)
        if self.file_manager:
            self.file_manager.top_bar.update_translations(lang_code)
            self.file_manager.sidebar.update_translations(lang_code)
            self.file_manager.action_bar.update_translations(lang_code)
            self.file_manager.statistics_view.update_translations(lang_code)
            self.file_manager.home_view.update_translations(lang_code)
            if hasattr(self.file_manager, 'sharing_view'):
                self.file_manager.sharing_view.update_translations(lang_code)
            if hasattr(self.file_manager, 'settings_overlay') and self.file_manager.settings_overlay:
                self.file_manager.settings_overlay.update_translations(lang_code)


def run_server():
    log = logging.getLogger('werkzeug')
    log.setLevel(logging.ERROR)
    flask_app.run(port=5000, use_reloader=False, debug=False)


if __name__ == "__main__":
    server_thread = threading.Thread(target=run_server, daemon=True)
    server_thread.start()

    app = QApplication(sys.argv)
    app.setStyle('Fusion')

    window = KemaslahApp()
    window.show()

    sys.exit(app.exec())